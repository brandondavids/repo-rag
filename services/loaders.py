from __future__ import annotations
import mailbox
from dataclasses import dataclass
from email import policy
from email.message import EmailMessage, Message
from email.parser import BytesParser
from pathlib import Path
import fitz
from bs4 import BeautifulSoup
from pptx import Presentation

SUPPORTED_EXTENSIONS = {".pdf", ".md", ".txt", ".eml", ".mbox", ".pptx"}

@dataclass(frozen=True)
class LoadedDocument:
    text: str
    metadata: dict[str, str | int | float | bool]

def load_file(path: Path) -> list[LoadedDocument]:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _load_pdf(path)
    if suffix in {".md", ".txt"}:
        return _load_text(path)
    if suffix == ".eml":
        return _load_eml(path)
    if suffix == ".mbox":
        return _load_mbox(path)
    if suffix == ".pptx":
        return _load_pptx(path)
    raise ValueError(f"Unsupported file type: {path.suffix}")

def _base_metadata(path: Path) -> dict[str, str]:
    return {
        "source": path.name,
        "source_path": str(path),
        "type": path.suffix.lower().lstrip("."),
    }

def _load_text(path: Path) -> list[LoadedDocument]:
    text = path.read_text(encoding="utf-8", errors="replace")
    return [LoadedDocument(text=text, metadata=_base_metadata(path))]

def _load_pdf(path: Path) -> list[LoadedDocument]:
    documents: list[LoadedDocument] = []

    with fitz.open(path) as pdf:
        for page_index, page in enumerate(pdf, start=1):
            text = page.get_text("text").strip()
            if not text:
                continue
            metadata = _base_metadata(path)
            metadata["page"] = page_index
            documents.append(LoadedDocument(text=text, metadata=metadata))
    
    return documents

def _extract_email_text(message: Message) -> str:
    plain_parts: list[str] = []
    html_parts: list[str] = []

    if message.is_multipart():
        for part in message.walk():
            content_disposition = part.get_content_disposition()
            if content_disposition == "attachment":
                continue
            content_type = part.get_content_type()
            payload = part.get_payload(decode=True)
            if payload is None:
                continue
            charset = part.get_content_charset() or "utf-8"
            decoded = payload.decode(charset, errors="replace")
            if content_type == "text/plain":
                plain_parts.append(decoded)
            elif content_type == "text/html":
                html_parts.append(decoded)
    else:
        payload = message.get_payload(decode=True)
        if payload is not None:
            charset = message.get_content_charset() or "utf-8"
            decoded = payload.decode(charset, errors="replace")
            if message.get_content_type() == "text/html":
                html_parts.append(decoded)
            else:
                plain_parts.append(decoded)
    
    if plain_parts:
        return "\n\n".join(plain_parts).strip()
    
    html = "\n\n".join(html_parts)
    return BeautifulSoup(html, "html.parser").get_text("\n").strip()

def _email_metadata(path: Path, message: Message, index: int | None = None) -> dict[str, str | int]:
    metadata: dict[str, str | int] = _base_metadata(path)
    metadata["subject"] = str(message.get("subject", ""))
    metadata["from"] = str(message.get("from", ""))
    metadata["date"] = str(message.get("date", ""))
    
    if index is not None:
        metadata["message_index"] = index
    
    return metadata

def _load_eml(path: Path) -> list[LoadedDocument]:
    with path.open("rb") as file:
        message: EmailMessage = BytesParser(policy=policy.default).parse(file)

    text = _extract_email_text(message)
    if not text:
        return []
    
    return [LoadedDocument(text=text, metadata=_email_metadata(path, message))]

def _load_mbox(path: Path) -> list[LoadedDocument]:
    documents: list[LoadedDocument] = []
    box = mailbox.mbox(path)
    for index, message in enumerate(box):
        text = _extract_email_text(message)
        if not text:
            continue
        documents.append(
            LoadedDocument(
                text=text,
                metadata=_email_metadata(path, message, index=index),
            )
        )
    
    return documents

def _load_pptx(path: Path) -> list[LoadedDocument]:
    documents: list[LoadedDocument] = []
    presentation = Presentation(str(path))
    for slide_index, slide in enumerate(presentation.slides, start=1):
        text_parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    paragraph_text = paragraph.text.strip()
                    if paragraph_text:
                        text_parts.append(paragraph_text)
            if shape.has_table:
                for row in shape.table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_texts:
                        text_parts.append(" | ".join(row_texts))
        slide_text = "\n".join(text_parts).strip()
        if not slide_text:
            continue
        metadata = _base_metadata(path)
        metadata["slide"] = slide_index
        documents.append(LoadedDocument(text=slide_text, metadata=metadata))
    
    return documents